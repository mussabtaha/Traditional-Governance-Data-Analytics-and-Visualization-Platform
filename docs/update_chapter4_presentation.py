from __future__ import annotations

import os
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
SHOTS = DOCS / "chapter4-screenshots"
SOURCE = DOCS / "Chapter_4_Results_and_Discussion.pptx"
OUTPUT = DOCS / "Chapter_4_Results_and_Discussion_Updated.pptx"

GREEN = "123524"
DEEP = "0B2D20"
GOLD = "C8A96A"
CREAM = "F8F7F3"
WHITE = "FFFFFF"
INK = "14251D"
MUTED = "65736B"
PALE = "E9EFEA"
LINE = "D7DED9"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


prs = Presentation(str(SOURCE))
original_slides = list(prs.slides)
original_slide_count = len(original_slides)


def remove_all_shapes(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def set_background(slide, color: str = CREAM) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    shape.name = "REV_BACKGROUND"


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.03,
    italic: bool = False,
    name: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_panel(slide, x: float, y: float, w: float, h: float, fill: str = WHITE, line: str = LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    return shape


def add_header(slide, section: str, title: str, subtitle: str, slide_number: int) -> None:
    set_background(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(GREEN); bar.line.fill.background()
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.28), prs.slide_width, Inches(0.22))
    footer.fill.solid(); footer.fill.fore_color.rgb = rgb(GREEN); footer.line.fill.background()
    add_text(slide, section.upper(), 0.55, 0.31, 4.8, 0.22, size=9, color=GREEN, bold=True, name="REV_SECTION")
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.63), Inches(0.46), Inches(0.025))
    rule.fill.solid(); rule.fill.fore_color.rgb = rgb(GOLD); rule.line.fill.background()
    add_text(slide, title, 0.55, 0.72, 12.05, 0.52, size=25, color=INK, font="Georgia", name="REV_TITLE")
    add_text(slide, subtitle, 0.57, 1.29, 11.95, 0.35, size=11.5, color=MUTED, name="REV_SUBTITLE")
    add_text(slide, "Chapter 4 • Results and Discussion", 0.55, 7.31, 5.0, 0.11, size=7.5, color="DDE7E0", margin=0)
    add_text(slide, str(slide_number), 12.55, 7.31, 0.3, 0.11, size=8, color=WHITE, align=PP_ALIGN.RIGHT, margin=0)


def image_contain(slide, filename: str, x: float, y: float, w: float, h: float):
    file_path = SHOTS / filename
    if not file_path.exists():
        raise FileNotFoundError(file_path)
    with Image.open(file_path) as image:
        width, height = image.size
    image_ratio = width / height
    box_ratio = w / h
    if image_ratio > box_ratio:
        draw_w = w
        draw_h = w / image_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * image_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    add_panel(slide, x - 0.035, y - 0.035, w + 0.07, h + 0.07, WHITE, LINE)
    return slide.shapes.add_picture(str(file_path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def add_caption(slide, text: str, x: float, y: float, w: float) -> None:
    add_text(slide, text, x, y, w, 0.22, size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER, margin=0, name="REV_CAPTION")


def add_explanation_box(slide, number: int, label: str, text: str, x: float, y: float, w: float, h: float = 0.86) -> None:
    add_panel(slide, x, y, w, h, WHITE, LINE)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.19), Inches(0.34), Inches(0.34))
    badge.fill.solid(); badge.fill.fore_color.rgb = rgb(GOLD); badge.line.color.rgb = rgb(WHITE)
    add_text(slide, str(number), x + 0.205, y + 0.263, 0.29, 0.13, size=8.5, color=DEEP, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, label.upper(), x + 0.64, y + 0.13, w - 0.82, 0.17, size=8, color=GREEN, bold=True, margin=0)
    add_text(slide, text, x + 0.64, y + 0.37, w - 0.82, h - 0.46, size=10.5, color=INK, bold=True, margin=0)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def prepare_slide(index: int, section: str, title: str, subtitle: str):
    slide = original_slides[index - 1]
    remove_all_shapes(slide)
    add_header(slide, section, title, subtitle, index)
    return slide


# Slide 2 — Introduction
slide = prepare_slide(2, "4.1 • Introduction", "What the final system demonstrates", "Chapter 4 presents the completed interface and the results users can obtain from it.")
for i, (label, text) in enumerate([
    ("EXPLORE", "Browse and inspect traditional governance records."),
    ("UNDERSTAND", "View clear summaries, maps, and descriptive charts."),
    ("COMPARE", "Compare countries, continents, and regions."),
    ("ANALYZE", "Review statistical relationships and their interpretation."),
]):
    x = 0.68 + (i % 2) * 6.16
    y = 1.92 + (i // 2) * 2.08
    add_panel(slide, x, y, 5.78, 1.75, WHITE, LINE)
    add_text(slide, f"0{i+1}", x + 0.28, y + 0.25, 0.7, 0.38, size=24, color=GOLD, bold=True, font="Georgia", margin=0)
    add_text(slide, label, x + 1.12, y + 0.27, 2.2, 0.22, size=10, color=GREEN, bold=True, margin=0)
    add_text(slide, text, x + 0.3, y + 0.9, 5.15, 0.45, size=14, color=INK, bold=True, font="Georgia", margin=0)
add_text(slide, "The following slides explain the final system from the user's perspective: what each page is, what it shows, what the user can do, and the benefit it provides.", 0.9, 6.25, 11.55, 0.55, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This slide introduces Chapter 4. It presents the final system and the results that a user can obtain from it. I will explain each page in simple terms: what it is, what it shows, what the user can do, and the benefit of the page.")


# Slide 3 — Final system overview from the user's perspective
slide = prepare_slide(3, "4.2 • Final System Overview", "A connected research journey", "The platform guides the user from a broad overview to detailed evidence and statistical interpretation.")
journey = [
    ("Home", "See the platform overview"),
    ("Groups", "Find relevant records"),
    ("Map & Statistics", "Understand patterns"),
    ("Comparison", "Compare geographic areas"),
    ("Analysis", "Interpret relationships"),
    ("About & Contact", "Understand and engage"),
]
for i, (label, text) in enumerate(journey):
    x = 0.44 + i * 2.12
    add_panel(slide, x, 2.15, 1.72, 1.48, PALE if i in (0, 4) else WHITE, GOLD if i == 4 else LINE)
    add_text(slide, str(i + 1), x + 0.63, 2.35, 0.45, 0.3, size=18, color=GOLD, bold=True, font="Georgia", align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, label, x + 0.13, 2.78, 1.46, 0.24, size=11.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, text, x + 0.13, 3.09, 1.46, 0.32, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, margin=0)
    if i < len(journey) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 1.79), Inches(2.67), Inches(0.26), Inches(0.42))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = rgb(GOLD); arrow.line.fill.background()
add_panel(slide, 1.05, 4.55, 11.25, 1.45, WHITE, LINE)
add_text(slide, "USER BENEFIT", 1.42, 4.88, 1.35, 0.2, size=9, color=GOLD, bold=True, margin=0)
add_text(slide, "The user can move from a global summary to filtered records, geographic comparisons, and statistical conclusions without leaving the same research platform.", 2.85, 4.76, 8.95, 0.55, size=15, color=INK, bold=True, font="Georgia", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This is an overview of the user's journey through the platform. The user begins on the Home page, explores records, studies maps and statistics, compares geographic areas, and then reviews statistical analysis. The About and Contact pages provide context and communication support.")


# Slide 4 — Home page
slide = prepare_slide(4, "4.3 • Home Page", "The starting point for exploration", "The Home page introduces the platform and presents the main dataset overview.")
image_contain(slide, "36-home-current.png", 0.58, 1.77, 8.25, 4.82)
add_caption(slide, "Figure 4.1: Home Page of the Traditional Governance Platform", 0.58, 6.65, 8.25)
for i, item in enumerate([
    ("WHAT IS IT?", "The main entry page of the platform."),
    ("WHAT DOES IT SHOW?", "Summary indicators, key tools, a world map, and recent records."),
    ("WHAT CAN THE USER DO?", "Open Groups, Map, Statistics, or Comparison."),
    ("WHAT IS THE BENEFIT?", "It gives a quick understanding of the dataset and where to begin."),
]):
    add_explanation_box(slide, i + 1, item[0], item[1], 9.08, 1.8 + i * 1.13, 3.52, 0.95)
set_notes(slide, "This is the Home page. It gives the user a quick overview of the platform and the dataset. It shows summary indicators, the main exploration tools, a geographic map, and recent group records. The user can use the buttons and navigation links to begin exploring the page that is most relevant to the research task.")


# Slide 5 — Home benefits
slide = prepare_slide(5, "4.3 • Home Page", "What the user gains from the dashboard", "The Home page converts a large dataset into an understandable first view.")
for i, (value, label, text) in enumerate([
    ("1,557", "GROUPS", "Shows the overall dataset size."),
    ("130", "COUNTRIES", "Shows the international coverage."),
    ("5", "CONTINENTS", "Introduces the geographic distribution."),
    ("1,351", "WITH TPI", "Shows how many records include a traditional institution."),
]):
    x = 0.65 + i * 3.08
    add_panel(slide, x, 1.95, 2.8, 1.72, WHITE, LINE)
    add_text(slide, value, x + 0.2, 2.22, 2.4, 0.4, size=25, color=GREEN, bold=True, font="Georgia", align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, label, x + 0.2, 2.77, 2.4, 0.2, size=9, color=GOLD, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, text, x + 0.28, 3.08, 2.24, 0.31, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, margin=0)
add_panel(slide, 0.85, 4.35, 11.63, 1.55, PALE, "C6D8CD")
add_text(slide, "HOW THE USER BEGINS", 1.22, 4.72, 2.05, 0.2, size=9, color=GREEN, bold=True, margin=0)
add_text(slide, "The user can choose a direct path: browse records, view geographic distribution, read descriptive statistics, or compare areas. This reduces the time needed to find the correct part of the system.", 3.18, 4.55, 8.72, 0.65, size=15, color=INK, bold=True, font="Georgia", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This slide explains the benefit of the Home page. The summary indicators tell the user how large and geographically broad the dataset is. The four main tools provide direct access to record exploration, geographic viewing, statistics, and comparison.")


# Slide 6 — Groups default and filtered
slide = prepare_slide(6, "4.4 • Groups Page", "Browse and narrow the traditional group records", "Search and filters help the user focus on records that match a research interest.")
image_contain(slide, "37-groups-default.png", 0.55, 1.72, 6.02, 4.18)
image_contain(slide, "38-groups-filtered-current.png", 6.77, 1.72, 6.02, 4.18)
add_caption(slide, "Figure 4.2: Default Group Records Interface", 0.55, 5.96, 6.02)
add_caption(slide, "Figure 4.3: Groups Filtered by Africa and the Search Term ‘Arab’", 6.77, 5.96, 6.02)
add_text(slide, "The user can search by group name and combine country, continent, region, leadership, and recognition filters. The visible records change to match the selected criteria.", 0.85, 6.42, 11.62, 0.45, size=12.5, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This is the Groups page. The left image shows the normal list of records. The right image shows a real example after searching for Arab and selecting Africa. The user can combine the available filters to narrow the results and find groups that are relevant to a country, region, leadership type, or recognition question.")


# Slide 7 — Group details
slide = prepare_slide(7, "4.4 • Groups Page", "Inspect a selected group in detail", "The detailed view organizes geographic, leadership, governance-function, and recognition information.")
image_contain(slide, "06-group-detail.png", 0.58, 1.72, 7.75, 4.94)
add_caption(slide, "Figure 4.4: Detailed Traditional Group Record", 0.58, 6.7, 7.75)
for i, item in enumerate([
    ("WHAT IS IT?", "A focused profile for one selected group."),
    ("WHAT DOES IT SHOW?", "Location, population, leadership, functions, and recognition."),
    ("WHAT CAN THE USER DO?", "Read the available fields without leaving the explorer."),
    ("WHAT IS THE BENEFIT?", "It supports record-level inspection after broader filtering."),
]):
    add_explanation_box(slide, i + 1, item[0], item[1], 8.65, 1.78 + i * 1.12, 3.9, 0.94)
set_notes(slide, "This is the detailed group view. After finding a relevant record, the user can open it and inspect the available information. The page clearly displays Yes, No, and Not Available values. The benefit is that the user can move from a broad list to detailed evidence for one group.")


# Slide 8 — Map
slide = prepare_slide(8, "4.5 • Interactive Map", "See how groups are distributed geographically", "The map displays the number of traditional groups represented in each continent category.")
image_contain(slide, "39-map-default-current.png", 0.58, 1.72, 8.28, 4.92)
add_caption(slide, "Figure 4.5: Interactive Geographic Map", 0.58, 6.68, 8.28)
for i, item in enumerate([
    ("WHAT IS IT?", "A visual geographic overview of the dataset."),
    ("WHAT DOES IT SHOW?", "A world map with a group total for each continent."),
    ("WHAT CAN THE USER DO?", "Select a marker to open the related group records."),
    ("WHAT IS THE BENEFIT?", "It makes geographic concentration easy to understand."),
]):
    add_explanation_box(slide, i + 1, item[0], item[1], 9.1, 1.78 + i * 1.12, 3.48, 0.94)
set_notes(slide, "This is the Interactive Map section. It represents the geographic distribution of the traditional groups. Each marker displays the number of records in a continent. The user can select a marker to continue to the related records. The benefit is a quick visual understanding of where the dataset is concentrated.")


# Slide 9 — Map interaction
slide = prepare_slide(9, "4.5 • Interactive Map", "A map selection opens the matching records", "The geographic view supports direct movement from a marker to a filtered list.")
image_contain(slide, "40-map-africa-result-current.png", 0.58, 1.72, 8.2, 4.92)
add_caption(slide, "Figure 4.6: Africa Map Marker Result", 0.58, 6.68, 8.2)
steps = [("1", "SELECT", "Africa marker"), ("2", "OPEN", "Groups page"), ("3", "FILTER", "Continent = Africa"), ("4", "LEARN", "726 matching records")]
for i, (number, label, text) in enumerate(steps):
    y = 1.92 + i * 1.08
    add_panel(slide, 9.12, y, 3.4, 0.82, PALE if i == 3 else WHITE, GOLD if i == 3 else LINE)
    add_text(slide, number, 9.35, y + 0.2, 0.35, 0.25, size=17, color=GOLD, bold=True, font="Georgia", align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, label, 9.85, y + 0.13, 0.8, 0.18, size=8.5, color=GREEN, bold=True, margin=0)
    add_text(slide, text, 9.85, y + 0.4, 2.25, 0.2, size=11.5, color=INK, bold=True, margin=0)
set_notes(slide, "This slide shows a real map interaction. The user selected Africa, and the Groups page opened with the Africa filter already applied. The result shows 726 matching records. This helps the user move naturally from a visual overview to the detailed data.")


# Slide 10 — Statistics
slide = prepare_slide(10, "4.6 • Statistics Page", "A visual summary of the dataset", "The page brings the most important descriptive patterns together in one dashboard.")
image_contain(slide, "41-statistics-current.png", 0.58, 1.72, 8.28, 4.92)
add_caption(slide, "Figure 4.7: Descriptive Statistics Dashboard", 0.58, 6.68, 8.28)
for i, item in enumerate([
    ("WHAT IS IT?", "The descriptive overview of the dataset."),
    ("WHAT DOES IT SHOW?", "Leadership, selection, functions, recognition, geography, population, and top countries."),
    ("WHAT CAN THE USER DO?", "Read charts and select a geographic scope."),
    ("WHAT IS THE BENEFIT?", "It makes broad dataset patterns easier to understand."),
]):
    add_explanation_box(slide, i + 1, item[0], item[1], 9.1, 1.78 + i * 1.12, 3.48, 0.94)
set_notes(slide, "This is the Statistics page. It provides a visual and descriptive summary of the dataset. The user can see leadership distribution, selection methods, governance functions, formal recognition, geographic coverage, largest groups, and top countries. The charts make the main patterns easier to understand.")


# Slide 11 — Scoped Statistics
slide = prepare_slide(11, "4.6 • Statistics Page", "Focus the dashboard on a geographic area", "The selected scope changes the summary indicators and every chart on the page.")
image_contain(slide, "42-statistics-africa-current.png", 0.58, 1.72, 8.55, 4.92)
add_caption(slide, "Figure 4.8: Africa-Scoped Statistics Dashboard", 0.58, 6.68, 8.55)
for i, (label, value) in enumerate([("GROUPS", "726"), ("COUNTRIES", "50"), ("WITH TPI", "685"), ("RECOGNIZED", "445")]):
    y = 1.85 + i * 1.03
    add_panel(slide, 9.45, y, 2.85, 0.8, WHITE, LINE)
    add_text(slide, label, 9.7, y + 0.14, 1.25, 0.17, size=8.5, color=GREEN, bold=True, margin=0)
    add_text(slide, value, 11.25, y + 0.13, 0.7, 0.26, size=18, color=GOLD, bold=True, font="Georgia", align=PP_ALIGN.RIGHT, margin=0)
add_text(slide, "The benefit is a focused description of one country, continent, or region instead of the whole dataset.", 9.52, 6.1, 2.75, 0.48, size=11.2, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This example shows the Statistics page after selecting Africa. All summary values and charts reflect the selected continent. The user can therefore study one geographic area instead of viewing only the global dataset.")


# Slide 12 — Comparison overview with continent evidence
slide = prepare_slide(12, "4.7 • Geographic Comparison", "Compare two geographic entities side by side", "The user selects a comparison level, Entity A, and Entity B before reviewing their differences.")
image_contain(slide, "44-comparison-continent-current.png", 0.58, 1.72, 8.35, 4.94)
add_caption(slide, "Figure 4.9: Continent Comparison between Africa and Asia", 0.58, 6.7, 8.35)
for i, item in enumerate([
    ("WHAT IS IT?", "A side-by-side geographic comparison page."),
    ("WHAT DOES IT SHOW?", "Group totals, leadership, recognition, functions, and population."),
    ("WHAT CAN THE USER DO?", "Choose Country, Continent, or Region and select two entities."),
    ("WHAT IS THE BENEFIT?", "Differences become easier to identify using aligned indicators."),
]):
    add_explanation_box(slide, i + 1, item[0], item[1], 9.2, 1.78 + i * 1.12, 3.35, 0.94)
set_notes(slide, "This is the Geographic Comparison page. The user first chooses the comparison type, then selects Entity A and Entity B, and starts the comparison. The page shows aligned indicators for leadership, recognition, governance functions, and population. This example compares Africa and Asia.")


# Slide 13 — Country and region evidence
slide = prepare_slide(13, "4.7 • Geographic Comparison", "Country, continent, and region are all supported", "The following completed examples demonstrate the remaining comparison levels.")
image_contain(slide, "43-comparison-country-current.png", 0.55, 1.72, 6.05, 4.55)
image_contain(slide, "45-comparison-region-current.png", 6.73, 1.72, 6.05, 4.55)
add_caption(slide, "Figure 4.10: Country Comparison between Kenya and Nigeria", 0.55, 6.32, 6.05)
add_caption(slide, "Figure 4.11: Region Comparison between Sub-Saharan Africa and East Asia & Pacific", 6.73, 6.32, 6.05)
add_text(slide, "Together with the continent example, these screens verify all three final comparison levels without using Group-versus-Group comparison.", 0.88, 6.69, 11.55, 0.35, size=11.8, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This slide provides visual evidence for the other two comparison levels. The left side compares Kenya and Nigeria as countries. The right side compares Sub-Saharan Africa and East Asia and Pacific as regions. Together with the previous continent example, the system supports Country, Continent, and Region comparison.")


# Slide 14 — fixed Statistical Analysis overview
slide = prepare_slide(14, "4.8 • Statistical Analysis Interface", "Where the user reads the statistical evidence", "Two readable views replace the previous long, compressed page screenshot.")
image_contain(slide, "46-analysis-interface-upper-current.png", 0.52, 1.72, 6.18, 4.52)
image_contain(slide, "47-analysis-interface-result-current.png", 6.82, 1.72, 6.0, 4.52)
add_caption(slide, "Figure 4.12: Research Question, Method, and Sample Summary", 0.52, 6.29, 6.18)
add_caption(slide, "Figure 4.13: Statistical Result, Interpretation, and Frequency Tables", 6.82, 6.29, 6.0)
add_text(slide, "1 Research question  •  2 Statistical method  •  3 Sample size  •  4 Statistic and p-value  •  5 Effect size  •  6 Interpretation and tables", 0.8, 6.66, 11.75, 0.32, size=10.7, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(slide, "This is the Statistical Analysis page. The left image shows the research question, selected method, total records, valid sample, missing values, and significance level. The right image shows the test statistic, p-value, degrees of freedom, effect size, sample size, interpretation, and the observed and expected tables. The user can use these elements to understand both the statistical result and its meaning.")


# Add About and Contact slides to the existing presentation, then place them after Slide 14.
about_slide = prs.slides.add_slide(prs.slide_layouts[-1])
add_header(about_slide, "4.9 • About Page", "Understand the purpose and scope of the project", "The About page explains why the platform exists and what research activities it supports.", 15)
image_contain(about_slide, "48-about-current.png", 0.58, 1.72, 8.25, 4.92)
add_caption(about_slide, "Figure 4.14: About Page of the Traditional Governance Platform", 0.58, 6.68, 8.25)
for i, item in enumerate([
    ("WHAT IS IT?", "The background and purpose page for the project."),
    ("WHAT DOES IT SHOW?", "The platform's purpose, scope, and research experience."),
    ("WHAT CAN THE USER DO?", "Learn how the system supports discovery, analysis, comparison, and current insights."),
    ("WHAT IS THE BENEFIT?", "It provides context before the user interprets the data."),
]):
    add_explanation_box(about_slide, i + 1, item[0], item[1], 9.08, 1.78 + i * 1.12, 3.5, 0.94)
set_notes(about_slide, "This is the About page. It explains the purpose and scope of the Traditional Governance platform. It shows that the system supports discovery, analysis, comparison, and communication of current insights. The benefit is that users understand the role of the project before interpreting its data.")

contact_slide = prs.slides.add_slide(prs.slide_layouts[-1])
add_header(contact_slide, "4.10 • Contact Page", "Communicate with the project team", "The Contact page provides a structured form for project feedback, dataset questions, and collaboration enquiries.", 16)
image_contain(contact_slide, "49-contact-current.png", 0.58, 1.72, 8.1, 4.92)
add_caption(contact_slide, "Figure 4.15: Contact Page and Project Enquiry Form", 0.58, 6.68, 8.1)
for i, item in enumerate([
    ("WHAT IS IT?", "The communication page for the project."),
    ("WHAT DOES IT SHOW?", "Project contact details and an enquiry form."),
    ("WHAT CAN THE USER DO?", "Enter a name, email, subject, and message, then select Send Message."),
    ("WHAT IS THE BENEFIT?", "It supports feedback, questions, and possible research collaboration."),
]):
    add_explanation_box(contact_slide, i + 1, item[0], item[1], 8.95, 1.78 + i * 1.12, 3.63, 0.94)
add_text(contact_slide, "Verification note: no message was sent during this review; delivery depends on the configured contact service.", 8.98, 6.15, 3.55, 0.5, size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
set_notes(contact_slide, "This is the Contact page. It provides project contact information and a form with Full name, Email address, Subject, and Message fields. The user can prepare feedback, a dataset question, or a collaboration enquiry and select Send Message. No message was sent during this review, so I do not claim that an email reached an inbox in the local environment.")

slide_ids = prs.slides._sldIdLst
about_id = slide_ids[-2]
contact_id = slide_ids[-1]
slide_ids.remove(about_id)
slide_ids.remove(contact_id)
slide_ids.insert(14, about_id)
slide_ids.insert(15, contact_id)


# Update the section labels of the preserved analytical slides without changing their results.
section_updates = {
    14: "STATISTICAL RESULTS • READING GUIDE",
    15: "STATISTICAL RESULTS • ANALYSIS 1",
    16: "STATISTICAL RESULTS • ANALYSIS 2",
    17: "STATISTICAL RESULTS • ANALYSIS 3",
    18: "STATISTICAL RESULTS • ANALYSIS 4",
    19: "STATISTICAL RESULTS • ANALYSIS 5",
    20: "STATISTICAL RESULTS • ANALYSIS 6",
    21: "STATISTICAL RESULTS • ANALYSIS 6",
    22: "STATISTICAL RESULTS • ANALYSIS 7",
    23: "STATISTICAL RESULTS • DYNAMIC ENGINE",
    24: "KEY FINDINGS",
    25: "KEY FINDINGS • EFFECT SIZE",
    26: "DISCUSSION",
    27: "LIMITATIONS • MISSING DATA",
    28: "SYSTEM ACCESSIBILITY",
    29: "SYSTEM RESULT SUMMARY",
    30: "CHAPTER SUMMARY",
}

for original_index, label in section_updates.items():
    slide = original_slides[original_index]
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip().upper()
        if re.match(r"^4\.(?:8|9|10|11|12)\s", text):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = label if run.text.strip() else run.text
            break


# Make the effect-size ranges exactly consistent with the project scale.
range_replacements = {
    "0.0–0.1": "0.00–0.09",
    "0.1–0.3": "0.10–0.29",
    "0.3–0.5": "0.30–0.49",
    "0.5–0.7": "0.50–0.69",
    "0.7–1.0": "0.70–1.00",
    "0.0-0.1": "0.00–0.09",
    "0.1-0.3": "0.10–0.29",
    "0.3-0.5": "0.30–0.49",
    "0.5-0.7": "0.50–0.69",
    "0.7-1.0": "0.70–1.00",
}

for slide in prs.slides:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old, new in range_replacements.items():
                        run.text = run.text.replace(old, new)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip().lower() in {"very weak", "weak", "moderate", "strong", "very strong"}:
                        cell.text = cell.text.strip().title()


# Strengthen the missing-data explanation on the preserved limitation slide.
missing_slide = original_slides[27]
add_panel(missing_slide, 0.78, 6.08, 7.55, 0.64, PALE, "C6D8CD")
add_text(
    missing_slide,
    "Missing values were not treated as ‘No’. Inferential analyses used valid complete cases for the variables required by each test.",
    1.0,
    6.25,
    7.1,
    0.28,
    size=10.2,
    color=INK,
    bold=True,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
)
set_notes(missing_slide, "This slide explains the main data limitation. Missing values are different from No. The analyses used valid complete cases for the variables required by each test. Healing has the largest amount of missing data, so healing-related findings should be interpreted with the greatest caution.")


# Reinforce the leadership-occurrence explanation without altering Analysis 5.
analysis5_slide = original_slides[19]
set_notes(analysis5_slide, "This analysis asks whether leadership distribution differs across continents. King, Chief, and Headman are separate binary indicators and may overlap within one group. Therefore, leadership occurrences are not the same as the number of unique groups. The result is statistically significant with a weak effect, and it describes association rather than causation.")


# Update the title date in the copied presentation.
for shape in original_slides[0].shapes:
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = run.text.replace("8 August 2026", "9 August 2026")


# Renumber all figure captions sequentially after the inserted pages.
figure_number = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if re.search(r"Figure 4\.\d+:", shape.text):
            figure_number += 1
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = re.sub(r"Figure 4\.\d+:", f"Figure 4.{figure_number}:", run.text)


prs.save(str(OUTPUT))
print({
    "source": str(SOURCE),
    "output": str(OUTPUT),
    "original_slides": original_slide_count,
    "final_slides": len(prs.slides),
    "figures": figure_number,
    "added_slides": 2,
})
